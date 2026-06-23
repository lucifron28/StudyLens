import fitz  # PyMuPDF
import io
import docx
from pptx import Presentation
from django.core.files.base import File

class ExtractionError(Exception):
    pass

def extract_pdf_text(file_obj: File) -> str:
    """Extracts text from a PDF file."""
    try:
        # Read the file content into memory
        file_bytes = file_obj.read()
        
        # Rewind the file pointer so it can be saved normally by Django
        file_obj.seek(0)
        
        # Open PDF from memory stream
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        
        text_parts = []
        for page in doc:
            text_parts.append(page.get_text())
            
        doc.close()
        
        # Clean up excessive newlines
        extracted = "\n\n".join(text_parts).strip()
        # Fallback to empty string if no text is extractable (e.g., scanned PDF without OCR)
        return extracted
    except Exception as e:
        raise ExtractionError(f"Failed to extract PDF text: {str(e)}")

def extract_docx_text(file_obj: File) -> str:
    """Extracts text from a DOCX file."""
    try:
        file_bytes = file_obj.read()
        file_obj.seek(0)
        
        doc = docx.Document(io.BytesIO(file_bytes))
        text_parts = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(text_parts)
    except Exception as e:
        raise ExtractionError(f"Failed to extract DOCX text: {str(e)}")

def extract_pptx_text(file_obj: File) -> str:
    """Extracts text from a PPTX file."""
    try:
        file_bytes = file_obj.read()
        file_obj.seek(0)
        
        prs = Presentation(io.BytesIO(file_bytes))
        text_parts = []
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                text_parts.append("\n".join(slide_text))
        return "\n\n---\n\n".join(text_parts)
    except Exception as e:
        raise ExtractionError(f"Failed to extract PPTX text: {str(e)}")
